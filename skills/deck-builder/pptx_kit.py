"""
pptx_kit — reusable helpers for building well-structured, on-brand .pptx decks
with python-pptx. Battle-tested on dark corporate decks; theme-agnostic.

Pairs with:
  - svg_icons.py   : embed crisp vector SVG icons (with PNG fallback) into shapes
  - render_check.py: faithful local PNG render to self-verify before delivery
  - references/design_tokens.md : palettes, style recipes, type/space scales (MIT, MiniMax)
  - references/structure.md     : deck structure principles (assertion titles, 1-idea-per-slide)

Design-token tables adapted from MiniMax-AI/skills (MIT). Dynamic-fit and
text-measurement ideas informed by GongRzhe/Office-PowerPoint-MCP-Server (MIT).
"""
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn, nsdecls
from PIL import ImageFont, ImageDraw, Image
import os

EMU_PER_IN = 914400
EMU_PER_PT = 12700

# ---- legibility floors (rule 4) — a deck NEVER ships text below these ----
# Body/caption copy must be readable at projection/print distance. If text only
# fits below the floor, CUT or restructure the copy — never shrink past it.
FLOOR_BODY_PT = 14      # body & captions: prefer 15-18, never below 14
FLOOR_KICKER_PT = 13    # eyebrow/kicker labels
ABS_FLOOR_PT = 12       # hard absolute — nothing legible ever drops under this

def IN(inches): return int(round(inches * EMU_PER_IN))
def C(hexstr): return RGBColor.from_string(hexstr.lstrip("#"))

# ---- text measurement (Arial proxy for Calibri; good enough for layout) ----
_FONT = "/System/Library/Fonts/Supplemental/Arial.ttf"
_FONT_B = "/System/Library/Fonts/Supplemental/Arial Bold.ttf"
_draw = ImageDraw.Draw(Image.new("RGB", (10, 10)))

def line_height(sz_pt, factor=1.30):
    return int(sz_pt * EMU_PER_PT * factor)

def n_lines(text, sz_pt, width_emu, bold=False):
    """How many lines `text` wraps to in a box `width_emu` wide at `sz_pt`.
    ⚠️ METRIC CAVEAT: measured with the BUNDLED font, which is usually WIDER than
    the deck's Calibri — so this can over-count (a title predicted at 2 lines may
    render on 1). It feeds fit/overflow/vbalance, so a wrong header line-count
    skews a balance anchor. For a critical anchor (title line count, a card that
    sits right at a void boundary) verify the real line count in a faithful render
    and correct the anchor; never trust the estimate blind near a tight boundary."""
    f = ImageFont.truetype(_FONT_B if bold else _FONT, int(sz_pt * 4))
    maxw = width_emu / EMU_PER_PT * 4
    words, lines, cur = text.split(), [], ""
    for w in words:
        t = (cur + " " + w).strip()
        if _draw.textlength(t, font=f) <= maxw or not cur:
            cur = t
        else:
            lines.append(cur); cur = w
    if cur: lines.append(cur)
    return max(1, len(lines))

def fit_all(texts, width_emu, height_emu, start=16, minimum=FLOOR_BODY_PT, bold=False):
    """SAFEGUARD against overflow in a ROW of equal cards: the single font size
    (pt) at which EVERY text fits the shared body box (= min of each text's
    fit_size). Apply this one size to all bodies → none overflows, all uniform.
    Always size card bodies through this against the real available height
    (card_bottom - body_top - bottom_pad). A card must contain its content.
    `minimum` defaults to the body floor: if the longest text only fits below it,
    fit_size returns the floor anyway → SHORTEN the copy or ENLARGE the card."""
    return min(fit_size(t, width_emu, height_emu, start, minimum, bold) for t in texts)

def fit_size(text, width_emu, height_emu, start=18, minimum=FLOOR_BODY_PT, bold=False):
    """Largest font size (pt) at which `text` fits within width x height. The
    fill-rule helper: prefer growing text to fill a box over leaving voids.
    `minimum` is the legibility floor (default body floor); the result is never
    below it. If even `minimum` overflows the box, the copy is too long for the
    box — cut it, don't override the floor."""
    sz = start
    while sz > minimum:
        if n_lines(text, sz, width_emu, bold) * line_height(sz) <= height_emu:
            return sz
        sz -= 1
    return minimum

def audit_text_sizes(prs, body_floor=FLOOR_BODY_PT, abs_floor=ABS_FLOOR_PT):
    """Scan every run in the deck and flag any below the legibility floor.
    Run this SILENTLY before handing a deck back (rule 4). Returns a list of
    dicts {slide, shape, text, size, level} where level is 'hard' (< abs_floor,
    illegible — must fix) or 'soft' (< body_floor, sub-floor body/caption).
    Titles/kickers are not auto-classified here — eyeball the 'soft' hits and
    only keep genuine kickers (>= FLOOR_KICKER_PT). Empty list == clean."""
    offenders = []
    for si, slide in enumerate(prs.slides, 1):
        for shp in slide.shapes:
            if not shp.has_text_frame:
                continue
            for par in shp.text_frame.paragraphs:
                for r in par.runs:
                    if not r.text.strip():
                        continue
                    sz = r.font.size
                    pt = sz.pt if sz is not None else None
                    if pt is None or pt >= body_floor:
                        continue
                    offenders.append({
                        "slide": si, "shape": shp.name,
                        "text": r.text.strip()[:60], "size": pt,
                        "level": "hard" if pt < abs_floor else "soft",
                    })
    return offenders

def _panel_fill(shp):
    from pptx.oxml.ns import qn as _qn
    for sp in shp._element.iter(_qn('a:solidFill')):
        c = sp.find(_qn('a:srgbClr'))
        if c is not None:
            return c.get('val')
    return None

def _text_need(shp):
    """Estimated rendered text height (EMU): per paragraph, wrapped lines x
    line-height at the paragraph's LARGEST run size + space_before + frame
    top/bottom margins. Returns (need_emu, had_text)."""
    tf = shp.text_frame
    need = (tf.margin_top or 0) + (tf.margin_bottom or 0)
    had = False
    for par in tf.paragraphs:
        runs = [r for r in par.runs if r.text]
        txt = "".join(r.text for r in runs)
        sb = 0
        try:
            if par.paragraph_format.space_before:
                sb = int(par.paragraph_format.space_before)
        except Exception:
            pass
        if not txt.strip():
            need += sb
            continue
        had = True
        szs = [r.font.size.pt for r in runs if r.font.size is not None]
        sz = max(szs) if szs else 18
        bold = any(bool(r.font.bold) for r in runs)
        need += n_lines(txt, sz, shp.width, bold=bold) * line_height(sz) + sb
    return need, had

def audit_overflow(prs, tol_in=0.06, slide_bottom_in=7.20):
    """Flag text that OVERFLOWS the bottom of its CONTAINING CARD (renders past
    the card edge / into the card below). The geometric audit_overlaps() catches
    shape-vs-shape collisions but NOT text spilling past its container — a
    multi-line paragraph or an appended line can exceed a fixed-height textbox and
    bleed out of the card. This estimates rendered text height (see _text_need),
    finds the smallest filled panel that contains the textbox's top, and flags
    when the text's bottom passes the panel's bottom (or the slide bottom when the
    text sits on the bare background). Returns {slide, shape, over_in, ctx, text};
    empty == clean. Run SILENTLY before handing back — overflow is a HARD defect
    (rule 4). STILL pair it with a faithful render at >=180 DPI zoomed into the
    densest cards: the estimator is approximate (font metrics, kerning), so treat
    it as a candidate flag, the render as the verdict."""
    out = []
    for si, slide in enumerate(prs.slides, 1):
        panels = [s for s in slide.shapes
                  if s.width and s.height and s.width > IN(1.4) and s.height > IN(0.7)
                  and _panel_fill(s) is not None
                  and not (s.has_text_frame and s.text_frame.text.strip())]
        for shp in slide.shapes:
            if not shp.has_text_frame or shp.width is None or shp.height is None:
                continue
            need, had = _text_need(shp)
            if not had:
                continue
            text_bottom = shp.top + need
            # smallest panel whose box contains the textbox's top-left-ish point
            cx, cy = shp.left + shp.width // 2, shp.top + IN(0.05)
            container = None
            for p in panels:
                if p.left <= cx <= p.left + p.width and p.top <= cy <= p.top + p.height:
                    if container is None or (p.width * p.height) < (container.width * container.height):
                        container = p
            if container is not None:
                limit = container.top + container.height
                ctx = "card"
            else:
                limit = IN(slide_bottom_in)
                ctx = "slide"
            if text_bottom > limit + int(tol_in * EMU_PER_IN):
                out.append({
                    "slide": si, "shape": shp.name, "ctx": ctx,
                    "over_in": round((text_bottom - limit) / EMU_PER_IN, 2),
                    "text": shp.text_frame.text.strip()[:55],
                })
    return out


def vmeasure(prs, slide_index, header_band_in=1.5):
    """Return (header_bottom, content_top, content_bottom, top_gap, bottom_gap) in
    EMU for one slide — the raw numbers behind the void-symmetry check. 'header'
    = shapes whose top is within header_band_in of the slide top (kicker/title);
    'content' = everything below. Visual bottoms use estimated text height
    (_text_need) so oversized empty textboxes don't lie about where ink ends.
    Full-bleed backgrounds (≥90% slide width, top≈0) are ignored."""
    sw, sh = prs.slide_width, prs.slide_height
    slide = prs.slides[slide_index]
    header_bottom = 0; tops = []; bottoms = []
    for shp in slide.shapes:
        if shp.top is None or shp.width is None or shp.height is None:
            continue
        if shp.width >= sw * 0.9 and shp.top < IN(0.3):
            continue
        if shp.has_text_frame and shp.text_frame.text.strip():
            need, _ = _text_need(shp); bot = shp.top + need
        else:
            bot = shp.top + shp.height
        if shp.top < IN(header_band_in):
            header_bottom = max(header_bottom, bot)
        else:
            tops.append(shp.top); bottoms.append(bot)
    if not tops:
        return (header_bottom, None, None, None, None)
    ct, cb = min(tops), max(bottoms)
    return (header_bottom, ct, cb, ct - header_bottom, sh - cb)

def audit_vbalance(prs, tol_in=0.18, header_band_in=1.5):
    """Flag VOID ASYMMETRY: the empty gap between the header and the content
    block must equal the gap between the content block and the slide bottom edge
    (a bigger void above the list than below it is a defect). For each slide
    this measures top_gap = first-content-top − header-bottom and bottom_gap =
    slide-bottom − last-content-bottom (text bottoms estimated, not the inflated
    textbox height), and flags |top_gap − bottom_gap| > tol. THE FIX: never
    hardcode a content band (CT/CB) that ignores the header — compute the block's
    real height and position it with `place_block(header_bottom, slide_bottom,
    block_h, 0.5)` so the two voids are equal by construction. Also: if both gaps
    are large, the block is too small → enlarge it (bigger type/nodes/pitch) to
    exploit the space. Heuristic (header vs content split) → confirm odd hits in a
    render; a slide that's intentionally NOT a single centred block may false-flag."""
    out = []
    for i in range(len(prs.slides)):
        hb, ct, cb, tg, bg = vmeasure(prs, i, header_band_in)
        if ct is None:
            continue
        if abs(tg - bg) > IN(tol_in):
            out.append({"slide": i + 1, "top_gap_in": round(tg / EMU_PER_IN, 2),
                        "bottom_gap_in": round(bg / EMU_PER_IN, 2),
                        "delta_in": round((tg - bg) / EMU_PER_IN, 2)})
    return out


def audit_overlaps(prs, frac=0.22, min_area_in2=0.06):
    """Autonomous render-defect catcher: geometrically detect shapes that COLLIDE.
    Run SILENTLY before handing a deck back, alongside audit_text_sizes + a faithful
    render. Returns dicts {slide, a, b, kind, overlap} where overlap = intersection
    area / smaller shape's area.

    Flags the high-signal, rarely-intentional cases:
      - 'hard': a PICTURE overlapping a text box or another picture (e.g. a logo
        sitting on top of copy — the classic bad render) — fix it.
      - 'soft': two text boxes overlapping (often a real collision, but the
        name+right-aligned-year-on-one-line pattern shares a bbox by design — eyeball).
    Text-on-card / picture-on-card is NOT flagged: a 'card' is a large filled
    autoshape with no text, which content legitimately sits on. bbox-based, so a
    soft hit can be a false positive (padded text boxes) — confirm in the render."""
    EMU2 = 914400.0 * 914400.0

    def kind_of(sh):
        try:
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                return 'pic'
        except Exception:
            return None
        has_text = sh.has_text_frame and sh.text_frame.text.strip()
        if has_text:
            return 'text'
        # filled, sizeable, text-less autoshape => a card/panel content sits on
        if sh.width and sh.height and sh.width * sh.height > IN(2.0) * IN(1.0):
            return 'card'
        return 'deco'

    out = []
    for si, slide in enumerate(prs.slides, 1):
        items = []
        for sh in slide.shapes:
            if sh.left is None or sh.width is None:
                continue
            k = kind_of(sh)
            if k in ('card', 'deco', None):
                continue
            items.append((sh, k))
        for i in range(len(items)):
            for j in range(i + 1, len(items)):
                (a, ka), (b, kb) = items[i], items[j]
                ix = max(0, min(a.left + a.width, b.left + b.width) - max(a.left, b.left))
                iy = max(0, min(a.top + a.height, b.top + b.height) - max(a.top, b.top))
                inter = ix * iy
                if inter <= 0:
                    continue
                sm = min(a.width * a.height, b.width * b.height)
                if sm / EMU2 < min_area_in2 or inter / sm < frac:
                    continue
                kinds = {ka, kb}
                level = 'soft' if kinds == {'text'} else 'hard'
                out.append({"slide": si, "a": a.name, "b": b.name,
                            "kind": f"{ka}|{kb}", "overlap": round(inter / sm, 2),
                            "level": level})
    return out

# ---- shapes ----
def _noline_noshadow(sp):
    sp.line.fill.background(); sp.shadow.inherit = False
    return sp

def rect(slide, l, t, w, h, fill):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(int(l)), Emu(int(t)), Emu(int(w)), Emu(int(h)))
    sp.fill.solid(); sp.fill.fore_color.rgb = C(fill)
    return _noline_noshadow(sp)

def rrect(slide, l, t, w, h, fill, radius=0.08):
    """radius is the python-pptx adjustment fraction (0..0.5). For a perfect
    pill set radius so corner == height/2 (use radius=0.5 on a thin bar)."""
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(int(l)), Emu(int(t)), Emu(int(w)), Emu(int(h)))
    sp.fill.solid(); sp.fill.fore_color.rgb = C(fill)
    try: sp.adjustments[0] = radius
    except Exception: pass
    return _noline_noshadow(sp)

def disc(slide, l, t, d, fill):
    sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(int(l)), Emu(int(t)), Emu(int(d)), Emu(int(d)))
    sp.fill.solid(); sp.fill.fore_color.rgb = C(fill)
    return _noline_noshadow(sp)

# ---- text ----
def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Emu(int(l)), Emu(int(t)), Emu(int(w)), Emu(int(h)))
    tf = tb.text_frame; tf.word_wrap = True
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    tf.vertical_anchor = anchor
    return tb, tf

def add_run(p, text, sz, color, bold=False, font="Calibri", spc=None):
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold; r.font.name = font
    r.font.color.rgb = C(color)
    return r

def line(slide, l, t, w, h, text, sz, color, bold=False, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, font="Calibri"):
    """Single-paragraph text box (the workhorse)."""
    tb, tf = textbox(slide, l, t, w, h, anchor)
    p = tf.paragraphs[0]; p.alignment = align
    add_run(p, text, sz, color, bold, font)
    return tb

def runs_line(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """One paragraph, multiple coloured runs. `runs` = [(text, sz, color, bold), ...]."""
    tb, tf = textbox(slide, l, t, w, h, anchor)
    p = tf.paragraphs[0]; p.alignment = align
    for text, sz, color, bold in runs:
        add_run(p, text, sz, color, bold)
    return tb

# ---- composite helpers ----
def kicker(slide, x, y, text, color="5DD5E8", bar=True, sz=13):
    """Small letterspaced section label with an accent bar (brand pattern)."""
    if bar:
        rect(slide, x, y - 90000, 300000, 52000, color)
    line(slide, x, y, 6000000, 280000, text, sz, color, bold=True)

def stack_fill(card_l, card_t, card_w, card_h, items, pad=IN(0.27), min_gap=60000):
    """THE FILL RULE. Distribute a vertical stack of items across a card with
    equal gaps (before / between / after) so the content fills the card instead
    of floating at the top. Returns y-positions; you draw at those y's.

    `items` = list of element heights (EMU), top to bottom.
    Returns list of top-y for each item.
    """
    inner = card_h - 2 * pad
    content = sum(items)
    gap = max(min_gap, (inner - content) // (len(items) + 1))
    ys, y = [], card_t + pad + gap
    for hgt in items:
        ys.append(y); y += hgt + gap
    return ys

def stack_top(card_top, heights, pad=IN(0.2), gap=IN(0.12)):
    """TOP-ALIGNED stack (for a ROW of equal-size cards). Returns y-positions
    using the SAME heights for every card so corresponding elements line up
    across cards (icons on one line, titles on one line, bodies starting on one
    line). The variation in text length then shows ONLY at the bottom — which is
    correct. For a row, compute each element's shared height as the MAX across
    all cards (e.g. title height = max title lines * line_height) and pass that.
    Contrast with stack_fill (centers within ONE card → breaks cross-card
    alignment; use only for a single standalone card, never a row)."""
    ys, y = [], card_top + pad
    for h in heights:
        ys.append(y); y += h + gap
    return ys

def place_block(band_top, band_bottom, block_h, bias=0.5):
    """Return the TOP y at which to place a block of height `block_h` inside the
    free vertical band [band_top, band_bottom] so the leftover space is BALANCED
    above and below — not dumped on one side. `band_top` is the bottom of the
    element above (header/intro), `band_bottom` the top of the element below
    (footer/testimonial). bias=0.5 centres the block; <0.5 sits it higher, >0.5
    lower. Use this to avoid the #1 layout smell: a big void on one side of a
    block while the other side is cramped (asymmetric whitespace reads as a bug)."""
    free = max(0, (band_bottom - band_top) - block_h)
    return int(band_top + free * bias)

def max_lines(texts, sz_pt, width_emu, bold=False):
    """Max wrapped-line count over a set of texts (to reserve shared title/body
    height across a row of cards)."""
    return max(n_lines(t, sz_pt, width_emu, bold) for t in texts)

def grid(n, container_l, container_w, gap, cols):
    """X positions + cell width for `cols` columns inside a container."""
    cw = (container_w - (cols - 1) * gap) // cols
    xs = [container_l + i * (cw + gap) for i in range(cols)]
    return xs, cw

# ---- icons (paired with svg_icons.embed) ----
def icon(slide, png_path, cx, cy, size, key=None):
    """Place a raster icon centered at (cx,cy). Name it svgicon_<key> so
    svg_icons.embed() can later attach the crisp vector SVG with PNG fallback."""
    pic = slide.shapes.add_picture(png_path, Emu(int(cx - size / 2)), Emu(int(cy - size / 2)),
                                   Emu(int(size)), Emu(int(size)))
    if key:
        pic.name = "svgicon_" + key
    return pic

def clear_shapes(slide, keep_pictures=False):
    """Remove all shapes (optionally keep pictures, e.g. a logo) before rebuild."""
    for sh in list(slide.shapes):
        if keep_pictures and sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            continue
        sh._element.getparent().remove(sh._element)

def remove_footers(prs):
    """Delete page-number / running-footer text boxes across all slides."""
    removed = 0
    for s in prs.slides:
        for sh in list(s.shapes):
            if sh.has_text_frame:
                t = sh.text_frame.text.strip()
                if (t.isdigit() and len(t) <= 3 and sh.top and sh.top > 5900000):
                    sh._element.getparent().remove(sh._element); removed += 1
    return removed

def set_slide_background(slide, image_path):
    """Set a slide's background to a full-bleed image — python-pptx has no
    high-level API for this, so we edit <p:cSld>/<p:bg> directly. Removes any
    existing <p:bg> AND any full-bleed background PICTURE shape first (some decks
    fake the bg with a stretched picture, which would otherwise hide the new bg),
    then inserts a stretched blipFill as the first child of <p:cSld>.

    Use this to give a section its own background (rule 7 — background rhythm).
    Pass a FLAT 1920x1080 PNG and keep the motif subtle so it recedes behind the
    content. If you switch a section to a LIGHT background, recolor the text that
    sits directly on it (eyebrow/title/subtitle → dark); text inside dark cards
    stays light. See references/backgrounds.md for moodboard-background recipes."""
    cSld = slide._element.find(qn('p:cSld'))
    old = cSld.find(qn('p:bg'))
    if old is not None:
        cSld.remove(old)
    for sh in list(slide.shapes):
        if (sh.shape_type == MSO_SHAPE_TYPE.PICTURE and sh.width and sh.height
                and sh.width > IN(12) and sh.height > IN(6.8)):
            sh._element.getparent().remove(sh._element)
    _, rId = slide.part.get_or_add_image_part(image_path)
    bg_xml = ('<p:bg %s><p:bgPr><a:blipFill><a:blip r:embed="%s"/>'
              '<a:stretch><a:fillRect/></a:stretch></a:blipFill>'
              '<a:effectLst/></p:bgPr></p:bg>') % (nsdecls('p', 'a', 'r'), rId)
    cSld.insert(0, parse_xml(bg_xml))


def add_slide_at(prs, layout, index):
    """Add a slide on `layout` and move it to 0-based `index` (python-pptx only
    appends). New slides inherit no per-slide background — call set_slide_background
    with a sibling section's bg (extract it from a sibling's <p:bg> blip)."""
    slide = prs.slides.add_slide(layout)
    lst = prs.slides._sldIdLst
    new = list(lst)[-1]
    lst.remove(new); lst.insert(index, new)
    return slide


def hyperlink_to_slide(shape, target_slide):
    """Make `shape` jump to `target_slide` on click — for clickable tables of
    contents / menus. Adds a slide relationship + a shape-level
    <a:hlinkClick action='ppaction://hlinksldjump'>. Attach to EVERY shape of a
    card (bg + each text box) so the whole card is clickable, not just the text."""
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    rId = shape.part.relate_to(target_slide.part, RT.SLIDE)
    cNvPr = shape._element.find(qn('p:nvSpPr')).find(qn('p:cNvPr'))
    hl = cNvPr.makeelement(qn('a:hlinkClick'),
                           {qn('r:id'): rId, 'action': 'ppaction://hlinksldjump'})
    cNvPr.insert(0, hl)


def hug_card_top(strip, card_l, card_t, card_w, card_h, card_radius_frac,
                 thickness_emu, samples=14):
    """Reshape a top-accent bar so its colour HUGS the card's rounded top corners
    (full card width, flush to the card top, top corners following the card's exact
    arc, square bottom) — the accent reads as the card's top edge "coloured in",
    stopping at the contour. Use this instead of a free-floating pill bar: a thin
    pill can't carry the card's corner radius (a roundRect radius caps at half the
    bar height) so its small corners overhang the card's larger rounded corners.

    `strip` = the existing accent shape (its solidFill/line are preserved; only the
    geometry + position/size are rewritten). `card_*` describe the parent card;
    `card_radius_frac` is the card roundRect's adj fraction (e.g. 0.06). Repositions
    the strip to (card_l, card_t, card_w, thickness_emu)."""
    import math
    R = min(card_radius_frac * min(card_w, card_h), card_w / 2)
    W, H = int(card_w), int(thickness_emu)

    def xL(y):                                  # card left contour at height y
        return R - math.sqrt(max(0.0, R * R - (R - y) ** 2))

    pts = [(round(R), 0), (round(W - R), 0)]                     # flat top edge
    for k in range(1, samples + 1):                             # right arc down
        y = H * k / samples
        pts.append((round(W - xL(y)), round(y)))
    pts.append((round(xL(H)), H))                                # bottom edge
    for k in range(samples - 1, -1, -1):                        # left arc up
        y = H * k / samples
        pts.append((round(xL(y)), round(y)))
    move = '<a:moveTo><a:pt x="%d" y="%d"/></a:moveTo>' % pts[0]
    lns = ''.join('<a:lnTo><a:pt x="%d" y="%d"/></a:lnTo>' % p for p in pts[1:])
    geom = ('<a:custGeom xmlns:a="%s"><a:avLst/><a:gdLst/><a:ahLst/><a:cxnLst/>'
            '<a:rect l="0" t="0" r="%d" b="%d"/><a:pathLst><a:path w="%d" h="%d">'
            '%s%s<a:close/></a:path></a:pathLst></a:custGeom>'
            ) % ('http://schemas.openxmlformats.org/drawingml/2006/main',
                 W, H, W, H, move, lns)
    strip.left, strip.top, strip.width, strip.height = int(card_l), int(card_t), W, H
    spPr = strip._element.spPr
    spPr.replace(spPr.find(qn('a:prstGeom')), parse_xml(geom))
