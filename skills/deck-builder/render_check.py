import sys, io, os
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image, ImageDraw, ImageFont

PPTX = sys.argv[1]
SLIDES = [int(x) for x in sys.argv[2].split(",")]
OUT = sys.argv[3] if len(sys.argv) > 3 else "/tmp"

p = Presentation(PPTX)
SW, SH = p.slide_width, p.slide_height
SCALE = 1600 / SW
def X(e): return int(round(e * SCALE))

FB = "/System/Library/Fonts/Supplemental/Arial Bold.ttf"
FR = "/System/Library/Fonts/Supplemental/Arial.ttf"
FI = "/System/Library/Fonts/Supplemental/Arial Italic.ttf"
def px_of(sz_pt): return max(7, int(round(sz_pt * 12700 * SCALE)))
def font(sz_pt, bold=False, ital=False):
    path = FB if bold else (FI if ital else FR)
    try: return ImageFont.truetype(path, px_of(sz_pt))
    except: return ImageFont.load_default()

def rgb(c):
    try: return (c[0], c[1], c[2])
    except: return None

def wrap(draw, text, fnt, maxw):
    words = text.split(" ")
    lines, cur = [], ""
    for w in words:
        t = (cur + " " + w).strip()
        if draw.textlength(t, font=fnt) <= maxw or not cur:
            cur = t
        else:
            lines.append(cur); cur = w
    if cur: lines.append(cur)
    return lines

def bg_media(slide):
    # find bg blip embed in the slide xml, map via rels
    from pptx.oxml.ns import qn
    el = slide._element
    blips = el.findall('.//' + qn('p:bg') + '//' + qn('a:blip'))
    if not blips: return None
    rid = blips[0].get(qn('r:embed'))
    if not rid: return None
    try:
        part = slide.part.related_part(rid)
        return Image.open(io.BytesIO(part.blob)).convert("RGB")
    except: return None

for sidx in SLIDES:
    s = p.slides[sidx-1]
    W, H = X(SW), X(SH)
    bg = bg_media(s)
    img = Image.new("RGB", (W, H), (10, 17, 36))
    if bg is not None:
        img.paste(bg.resize((W, H)), (0, 0))
    d = ImageDraw.Draw(img, "RGBA")
    for sh in s.shapes:
        try: x, y, w, h = X(sh.left), X(sh.top), X(sh.width), X(sh.height)
        except: continue
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                im = Image.open(io.BytesIO(sh.image.blob)).convert("RGBA").resize((max(1,w), max(1,h)))
                img.paste(im, (x, y), im)
            except: d.rectangle([x,y,x+w,y+h], outline=(255,0,0))
            continue
        # fill
        fc = None
        try:
            if sh.fill.type == 1: fc = rgb(sh.fill.fore_color.rgb)
        except: pass
        if fc:
            rad = max(2, int(h*0.10))
            d.rounded_rectangle([x,y,x+w,y+h], radius=rad, fill=fc+(255,))
        # text
        if sh.has_text_frame and sh.text_frame.text.strip():
            tf = sh.text_frame
            va = str(getattr(tf, "vertical_anchor", None) or "")
            # gather lines first
            blocks = []
            for para in tf.paragraphs:
                r = para.runs[0] if para.runs else None
                sz = 12; col=(255,255,255); bold=False; ital=False
                if r is not None:
                    if r.font.size: sz = r.font.size.pt
                    if r.font.bold: bold = True
                    if r.font.italic: ital = True
                    try:
                        if r.font.color and r.font.color.rgb: col = rgb(r.font.color.rgb)
                    except: pass
                fnt = font(sz, bold, ital)
                al = getattr(para.alignment, "value", None)
                txt = para.text
                if not txt.strip():
                    blocks.append((("",), fnt, col, al, sz)); continue
                lines = wrap(d, txt, fnt, max(10, w-6))
                blocks.append((lines, fnt, col, al, sz))
            lh = lambda sz: int(px_of(sz)*1.32)+2
            total = sum(lh(sz)*len(lines) for lines,_,_,_,sz in blocks)
            if "BOTTOM" in va: ty = y + h - total
            elif "MIDDLE" in va or "CENTER" in va: ty = y + (h-total)//2
            else: ty = y + 2
            for lines, fnt, col, al, sz in blocks:
                for ln in lines:
                    tw = d.textlength(ln, font=fnt)
                    if al == 2: tx = x + (w-tw)//2
                    elif al == 3: tx = x + w - tw
                    else: tx = x + 4
                    if ln: d.text((tx, ty), ln, fill=col, font=fnt)
                    ty += lh(sz)
    out = os.path.join(OUT, f"r_{sidx}.png")
    img.save(out)
    print("rendered", out)
